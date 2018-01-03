from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm #Inches
from pptx.enum.chart import XL_LEGEND_POSITION

if __name__ == '__main__':
    prs = Presentation('template.pptx')
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = '报告'

    name_objects = ["object1", "object2", "object3"]
    name_AIs = ["AI1", "AI2", "AI3"]
    val_AI1 = (19.2, 21.4, 16.7)
    val_AI2 = (22.3, 28.6, 15.2)
    val_AI3 = (20.4, 26.3, 14.2)
    val_AIs = [val_AI1, val_AI2, val_AI3]

    rows = 4
    cols = 4
    top    = Cm(12.5)
    left   = Cm(3.5) #Inches(2.0)
    width  = Cm(24) # Inches(6.0)
    height = Cm(6) # Inches(0.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Cm(6)# Inches(2.0)
    table.columns[1].width = Cm(6)
    table.columns[2].width = Cm(6)
    table.columns[3].width = Cm(6)

    # write column headings
    table.cell(0, 1).text = name_objects[0]
    table.cell(0, 2).text = name_objects[1]
    table.cell(0, 3).text = name_objects[2]

    # write body cells
    table.cell(1, 0).text = name_AIs[0]
    table.cell(1, 1).text = str(val_AI1[0])
    table.cell(1, 2).text = str(val_AI1[1])
    table.cell(1, 3).text = str(val_AI1[2])

    table.cell(2, 0).text = name_AIs[1]
    table.cell(2, 1).text = str(val_AI2[0])
    table.cell(2, 2).text = str(val_AI2[1])
    table.cell(2, 3).text = str(val_AI2[2])

    table.cell(3, 0).text = name_AIs[2]
    table.cell(3, 1).text = str(val_AI3[0])
    table.cell(3, 2).text = str(val_AI3[1])
    table.cell(3, 3).text = str(val_AI3[2])

    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = name_objects
    chart_data.add_series(name_AIs[0], val_AI1)
    chart_data.add_series(name_AIs[1], val_AI2)
    chart_data.add_series(name_AIs[2], val_AI3)

    # add chart to slide --------------------
    x, y, cx, cy = Cm(3.5), Cm(4.2), Cm(24), Cm(8)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

    chart = graphic_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False

    value_axis = chart.value_axis
    value_axis.maximum_scale = 100.0

    value_axis.has_title = True
    value_axis.axis_title.has_text_frame = True
    value_axis.axis_title.text_frame.text = "False positive"
    value_axis.axis_title.text_frame.auto_size
    #fit_text(font_family='Calibri', max_size=18, bold=False, italic=False, font_file=None)

    prs.save('test_template.pptx')
